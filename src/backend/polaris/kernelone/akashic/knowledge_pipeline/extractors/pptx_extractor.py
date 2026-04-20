"""PowerPoint PPTX Extractor for Knowledge Pipeline.

Uses python-pptx for text extraction from .pptx files.
Provides graceful degradation if python-pptx is not installed.
"""

from __future__ import annotations

import logging
from typing import TYPE_CHECKING

from polaris.kernelone.akashic.knowledge_pipeline.extractors.base import (
    BaseExtractor,
)

logger = logging.getLogger(__name__)

# Try to import python-pptx; if not available, extractor will degrade gracefully
try:
    import pptx

    PYPPTX_AVAILABLE = True
except ImportError:
    PYPPTX_AVAILABLE = False
    pptx = None  # type: ignore[assignment]

if TYPE_CHECKING:
    from polaris.kernelone.akashic.knowledge_pipeline.extractors.base import (
        ExtractionOptions,
    )
    from polaris.kernelone.akashic.knowledge_pipeline.protocols import (
        DocumentInput,
        ExtractedFragment,
    )


class PptxExtractor(BaseExtractor):
    """Extractor for Microsoft PowerPoint .pptx files.

    Extracts text from slides, preserving slide structure.

    Supported MIME types:
    - application/vnd.openxmlformats-officedocument.presentationml.presentation
    """

    SUPPORTED_MIME_TYPES = ("application/vnd.openxmlformats-officedocument.presentationml.presentation",)

    def is_available(self) -> bool:
        """Check if PPTX extraction is available (python-pptx installed)."""
        return PYPPTX_AVAILABLE

    async def extract(self, doc: DocumentInput) -> list[ExtractedFragment]:  # type: ignore[override]
        """Extract text fragments from a PowerPoint document."""
        from polaris.kernelone.akashic.knowledge_pipeline.protocols import (
            ExtractedFragment,
        )

        if not PYPPTX_AVAILABLE:
            logger.warning(
                "PPTX extraction unavailable: python-pptx not installed. Install with: pip install python-pptx"
            )
            return self._fallback_extract(doc)

        if isinstance(doc.content, str):
            logger.warning("PPTX content provided as string, not bytes")
            return [
                ExtractedFragment(
                    text=doc.content[:5000],
                    line_start=1,
                    line_end=1,
                    mime_type=doc.mime_type,
                    metadata={"warning": "string_content_not_processed"},
                )
            ]

        try:
            import asyncio
            import io

            loop = asyncio.get_running_loop()

            def _parse_pptx_bytes() -> list[ExtractedFragment]:
                from pptx import Presentation

                assert not isinstance(doc.content, str)  # mypy: str | bytes → bytes
                prs = Presentation(io.BytesIO(doc.content))
                fragments: list[ExtractedFragment] = []
                global_line = 1

                for slide_num, slide in enumerate(prs.slides, start=1):
                    # Slide header
                    fragments.append(
                        ExtractedFragment(
                            text=f"[SLIDE {slide_num}]",
                            line_start=global_line,
                            line_end=global_line,
                            mime_type=doc.mime_type,
                            metadata={
                                "slide": slide_num,
                                "fragment_type": "pptx_slide_header",
                            },
                        )
                    )
                    global_line += 1

                    # Extract shapes
                    for shape in slide.shapes:
                        if not hasattr(shape, "text"):
                            continue
                        text = shape.text.strip()
                        if not text:
                            continue
                        fragments.append(
                            ExtractedFragment(
                                text=text,
                                line_start=global_line,
                                line_end=global_line,
                                mime_type=doc.mime_type,
                                metadata={
                                    "slide": slide_num,
                                    "fragment_type": "pptx_shape",
                                },
                            )
                        )
                        global_line += 1

                logger.debug(
                    "Extracted %d fragments from PPTX (%d slides)",
                    len(fragments),
                    len(prs.slides),
                )
                return fragments

            return await loop.run_in_executor(None, _parse_pptx_bytes)

        except (RuntimeError, ValueError) as exc:
            logger.warning("PPTX extraction failed: %s", exc)
            return self._fallback_extract(doc)

    def _fallback_extract(self, doc: DocumentInput) -> list[ExtractedFragment]:
        """Fallback when python-pptx is unavailable."""
        from polaris.kernelone.akashic.knowledge_pipeline.protocols import (
            ExtractedFragment,
        )

        if isinstance(doc.content, str):
            return [
                ExtractedFragment(
                    text=doc.content[:5000],
                    line_start=1,
                    line_end=1,
                    mime_type=doc.mime_type,
                    metadata={"fallback": "string_content"},
                )
            ]

        return [
            ExtractedFragment(
                text="[PPTX content could not be extracted - python-pptx not installed]",
                line_start=1,
                line_end=1,
                mime_type=doc.mime_type,
                metadata={"error": "python-pptx_missing"},
            )
        ]

    def _do_extract(  # type: ignore[override]
        self,
        text: str,
        options: ExtractionOptions,
    ) -> list[ExtractedFragment]:
        """Not used - override extract() handles bytes directly."""
        return []


__all__ = ["PptxExtractor"]
